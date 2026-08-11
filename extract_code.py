#!/usr/bin/env python3
"""
PMPP 강의 덱(pptx)에서 CUDA 코드 조각을 추출한다.

usage:  python3 extract_code.py <슬라이드_디렉터리> <출력_디렉터리>
출력:   <출력>/ch10_s16_reduction_code_with_shared_memory.cu  형태
        + index.md  (슬라이드 → 파일 매핑표)
"""
import re, sys, unicodedata
from pathlib import Path
from pptx import Presentation

# 코드로 판정할 신호
CODE_HINTS = (
    "__global__", "__shared__", "__syncthreads", "threadIdx", "blockIdx",
    "blockDim", "gridDim", "cudaMalloc", "cudaMemcpy", "atomicAdd",
    "for(", "for (", "if(", "if (", "};", "int ", "float ",
)


def looks_like_code(text: str) -> bool:
    lines = [l for l in text.splitlines() if l.strip()]
    if len(lines) < 2:
        return False
    hits = sum(1 for l in lines if any(h in l for h in CODE_HINTS))
    return hits >= max(2, len(lines) // 4)


def slug(s: str) -> str:
    s = unicodedata.normalize("NFKD", s)
    s = re.sub(r"[^\w\s-]", "", s).strip().lower()
    return re.sub(r"[\s_-]+", "_", s)[:48] or "untitled"


def main(src: Path, out: Path) -> None:
    out.mkdir(parents=True, exist_ok=True)
    rows = []

    for pptx in sorted(src.glob("*.pptx")):
        m = re.search(r"Chapter\s*(\d+)", pptx.name)
        ch = f"ch{int(m.group(1)):02d}" if m else slug(pptx.stem)[:6]

        for n, slide in enumerate(Presentation(pptx).slides, 1):
            texts = [
                sh.text_frame.text
                for sh in slide.shapes
                if sh.has_text_frame and sh.text_frame.text.strip()
            ]
            if not texts:
                continue
            title = texts[0].splitlines()[0].strip()

            for body in texts[1:]:
                if not looks_like_code(body):
                    continue
                name = f"{ch}_s{n:02d}_{slug(title)}.cu"
                header = (
                    f"// 출처: {pptx.name} / slide {n}\n"
                    f"// 제목: {title}\n"
                    f"// 주의: 커널 본체 조각이다. 시그니처와 호스트 코드는 직접 붙여야 한다.\n\n"
                )
                (out / name).write_text(header + body.rstrip() + "\n", encoding="utf-8")
                rows.append((ch, n, title, name, len(body.splitlines())))

    lines = ["# 추출된 코드 조각", "", "| 장 | 슬라이드 | 제목 | 파일 | 행 |", "|---|---|---|---|---|"]
    lines += [f"| {c} | {n} | {t} | `{f}` | {L} |" for c, n, t, f, L in rows]
    (out / "index.md").write_text("\n".join(lines) + "\n", encoding="utf-8")
    print(f"{len(rows)}개 조각 추출 → {out}")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        sys.exit(__doc__)
    main(Path(sys.argv[1]), Path(sys.argv[2]))
